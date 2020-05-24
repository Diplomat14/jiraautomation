from jiraautomation.operations.operation import basic_operation
from jiraautomation.operations.generate_wbs import generate_wbs
from pptx.shapes.graphfrm import GraphicFrame
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.slide import Slide
import os


class generate_presentation(basic_operation):

    @staticmethod
    def name():
        return "GeneratePresentation"

    @staticmethod
    def init_arguments(operation_group):
        operation_group.add_argument('-genpr_Layout', '--generatepr_Layout', required=False,
                                     help='Number of layout to use')
        operation_group.add_argument('-genpr_FontSize', '--generatepr_FontSize', required=False,
                                     help='Size of font')
        pass

    @staticmethod
    def parse_arguments(args):
        pass

    def __init__(self, iLogger, filename):
        super().__init__(iLogger)
        self.__filename = filename
        self.__presentation = None
        self.set_presentation()

    def execute(self, container, args):
        l = self.logger

        try:

            try:
                if self.__presentation:
                    self.__presentation.save(self.__filename)
                    l.msg('Presentation {} was successfully saved'.format(os.path.basename(self.__filename)))

            except Exception as e:
                l.error("Exception happened boards search " + str(e))

        except Exception as e:
            l.error("Exception happened during connection establishment " + str(e))

    def set_presentation(self):
        self.__presentation = Presentation(self.__filename)

    def get_slide_layout(self, layout):
        return self.__presentation.slide_layouts[layout]

    def get_slide(self, slide_number, layout=5):
        elem = self.get_slide_id(slide_number)
        if not elem:
            slide_layout = self.get_slide_layout(layout)
            self.logger.msg('Creating new slide with layout №{}'.format(layout))
            elem = self.__presentation.slides.add_slide(slide_layout)

        slide = SlideExt(elem.part._element, elem.part, self.logger)

        return slide

    def get_slide_id(self, number):
        try:
            return self.__presentation.slides[number - 1]
        except Exception:
            self.logger.msg('Slide №{} is not exist'.format(number))

    @property
    def slides_height(self):
        return self.__presentation.slide_height.cm


class SlideExt(Slide):

    def __init__(self, elem, part, iLogger):
        super().__init__(elem, part)
        self.logger = iLogger

    @property
    def slide_title(self):
        if self.shapes.title:
            return self.shapes.title.text

    @slide_title.setter
    def slide_title(self, value):
        if self.shapes.title:
            self.shapes.title.text = value

    def get_placeholder(self, name):
        if self.placeholders:
            for pls in self.placeholders:
                if pls.name == name:
                    return pls
            else:
                raise Exception('No placeholder with name "{}"'.format(name))
        else:
            raise Exception('No placeholders in slide')

    def change_placeholder_text(self, name, value):
        pls = self.get_placeholder(name)
        self.logger.msg('Adding data to placeholder "{}"'.format(name))
        pls.text = value

    @property
    def has_table(self):
        for elem in self.shapes:
            if isinstance(elem, GraphicFrame) and elem.table:
                return True

        return False

    @property
    def tables(self):
        tbls = [elem for elem in self.shapes if isinstance(elem, GraphicFrame) and elem.table]
        return tbls

    def get_table_by_name(self, name):
        if not self.has_table:
            raise Exception('No table on slide')

        for tbl in self.tables:
            if tbl.name == name:
                return tbl
        else:
            raise Exception('No table with name "{}"'.format(name))

    def modify_existing_table(self, tbl, rows, cols):
        top, left, width, height = SlideExt.get_table_params(tbl)
        SlideExt.remove_table(tbl)

        tbl = self.add_new_table(rows, cols, top, left, width, height)
        return tbl

    def add_new_table(self, rows, cols, top, left, width, height):
        top = Inches(top)
        left = Inches(left)
        width = Inches(width)
        height = Inches(height)

        self.logger.msg('Adding new table to slide')
        tbl = self.shapes.add_table(rows + 1, cols, left, top, width, height)
        return tbl

    # TODO: Check if table size is acceptable
    # def fit_table_size(self, top, rows, height, slides_height):
    #    size = (top + (rows + 1) * height)
    #    if size >= Inches(slides_height):
    #        if height > Inches(1):
    #            height -= Inches(1)
    #            return self.fit_table_size(top, rows, height)
    #        else:
    #            raise Exception('Too big table')
    #    else:
    #        return height

    @staticmethod
    def remove_table(table):
        elem = table.element
        elem.getparent().remove(elem)

    @staticmethod
    def get_table_params(table):
        return table.top.inches, table.left.inches, \
               table.width.inches, table.height.inches

    @staticmethod
    def add_data_to_table(shp, df, size=11):
        rows, cols = df.shape
        if cols > len(shp.table.columns) or rows > len(shp.table.rows):
            raise Exception('Data dimension > table dimension ({}x{} > {}x{})'.format(rows, cols, len(shp.table.rows),
                                                                                      len(shp.table.columns)))
        column_names = list(df.columns)

        for col_index, col_name in enumerate(column_names):
            shp.table.cell(0, col_index).text = col_name
            run = shp.table.cell(0, col_index).text_frame.paragraphs[0]
            run.font.size = Pt(size)

        data = df.values

        for row in range(rows):
            for col in range(cols):
                val = data[row, col]

                shp.table.cell(row + 1, col).text = str(val)
                run = shp.table.cell(row + 1, col).text_frame.paragraphs[0]
                run.font.size = Pt(size)
